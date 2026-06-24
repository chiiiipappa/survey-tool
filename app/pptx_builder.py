"""STEP4 レポートページ → PowerPoint (.pptx) 変換エンジン。
グラフ・集計表ともに PowerPoint ネイティブオブジェクトとして出力する。
"""

from __future__ import annotations

import io
from copy import deepcopy
from lxml import etree

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn


# ---------------------------------------------------------------------------
# 定数
# ---------------------------------------------------------------------------

CHART_COLORS = [
    "#0071BC", "#DF0515", "#3DAA68", "#F5A623", "#9B59B6",
    "#1ABC9C", "#E67E22", "#E74C3C", "#2980B9", "#27AE60",
]

CHART_TYPE_MAP: dict[str, XL_CHART_TYPE] = {
    "vbar":               XL_CHART_TYPE.COLUMN_CLUSTERED,
    "hbar":               XL_CHART_TYPE.BAR_CLUSTERED,
    "grouped_vbar":       XL_CHART_TYPE.COLUMN_CLUSTERED,
    "grouped_hbar":       XL_CHART_TYPE.BAR_CLUSTERED,
    "stacked100_vbar":    XL_CHART_TYPE.COLUMN_STACKED_100,
    "stacked100_hbar":    XL_CHART_TYPE.BAR_STACKED_100,
    "brand_hbar":         XL_CHART_TYPE.COLUMN_CLUSTERED,   # indexAxis="x" → 縦棒
    "brand_vbar":         XL_CHART_TYPE.BAR_CLUSTERED,      # indexAxis="y" → 横棒
    "brand_vbar_stacked": XL_CHART_TYPE.BAR_STACKED_100,
    "auto":               XL_CHART_TYPE.COLUMN_CLUSTERED,
    "small_multiples":    XL_CHART_TYPE.COLUMN_CLUSTERED,
}

# スライドサイズ: 16:9 widescreen
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# コンテンツ領域
MARGIN_L = Inches(0.47)
MARGIN_T_TITLE = Inches(0.08)
TITLE_H = Inches(0.47)
SUBTITLE_H = Inches(0.36)
CONTENT_TOP = Inches(1.0)
CONTENT_H = Inches(6.25)
CONTENT_W = SLIDE_W - MARGIN_L * 2

# フッター（n数・分析軸）
FOOTER_H             = Inches(0.35)
FOOTER_MARGIN_BOTTOM = Inches(0.07)

# 集計表罫線
TABLE_BORDER_PT = 0.25
TABLE_BORDER_COLOR = "676767"

FONT_NAME = "Noto Sans JP"
FONT_NAME_FALLBACK = "Noto Sans"


# ---------------------------------------------------------------------------
# パブリック API
# ---------------------------------------------------------------------------

def _count_page_split_charts(cc: dict) -> int:
    """1ページ分の分割グラフ件数を返す（splitDatasetIndices 優先、旧形式 fallback）。"""
    if cc.get("splitMode", "normal") not in ("by_axis", "by_comparison"):
        return 0
    indices = cc.get("splitDatasetIndices")
    if indices is not None:
        return len(indices)
    s = int(cc.get("splitChunkStart") or 0)
    e = int(cc.get("splitChunkEnd")   or 0)
    return max(0, e - s)


def count_split_charts(pages: list[dict]) -> int:
    """pages 内の分割グラフ件数合計を返す（検証用）。"""
    return sum(_count_page_split_charts(p.get("chartConfig") or {}) for p in pages)


def build_pptx(pages: list[dict], chart_results: list[dict]) -> tuple[int, int, bytes]:
    """全ページを 1 ファイルの PPTX に変換して bytes を返す。"""
    import logging
    logger = logging.getLogger(__name__)

    cr_map: dict[str, dict] = {cr["id"]: cr for cr in chart_results}

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # 完全空白レイアウト

    # --- 事前チェック: 分割グラフの件数集計 ---
    expected_split_total = count_split_charts(pages)
    actual_split_total = 0

    for page in pages:
        cr_id = (page.get("aggregationConfig") or {}).get("chartResultId", "")
        cr = cr_map.get(cr_id)
        if not cr:
            logger.warning(f"[PPTX] ChartResult 未発見: cr_id={cr_id}")
            continue
        slide = prs.slides.add_slide(blank_layout)
        cc = (page.get("chartConfig") or {})
        count = _count_page_split_charts(cc)
        if count > 0:
            actual_split_total += count
            logger.info(f"[PPTX] スライド分割グラフ: {count}件")
        _add_slide(prs, slide, page, cr)

        # tablePosition == "separate" → 集計表を別スライドに追加
        lc = page.get("layoutConfig") or {}
        if lc.get("tablePosition") == "separate":
            table_slide = prs.slides.add_slide(blank_layout)
            _add_table_only_slide(prs, table_slide, page, cr)

    # --- 事後検証 ---
    if expected_split_total > 0:
        logger.info(f"[PPTX] 分割グラフ合計: 期待={expected_split_total}件 / 出力={actual_split_total}件")
        if actual_split_total != expected_split_total:
            logger.error(
                f"[PPTX] ⚠️ 件数不一致！ 期待={expected_split_total}件 出力={actual_split_total}件"
            )

    buf = io.BytesIO()
    prs.save(buf)
    return expected_split_total, actual_split_total, buf.getvalue()


# ---------------------------------------------------------------------------
# スライド構築
# ---------------------------------------------------------------------------

def _build_footer_text(page: dict, cr: dict) -> str:
    """n数・分析軸のフッターテキストを生成する（フロントエンド _buildNCountHtml に対応）。"""
    parts: list[str] = []

    comp_ds     = cr.get("comparison_datasets") or []
    axis_cats   = cr.get("axis_categories") or []
    axis_totals = cr.get("axis_totals") or []

    if comp_ds:
        n_parts = [
            f"{ds.get('target_value', '')}: n={sum(ds.get('axis_totals') or [])}"
            for ds in comp_ds
        ]
        parts.append("  /  ".join(n_parts))
    elif axis_totals:
        n_parts: list[str] = []
        for i, n in enumerate(axis_totals):
            cat = axis_cats[i] if i < len(axis_cats) else ""
            n_parts.append(f"n={n}" if cat in ("全体", "") else f"{cat}: n={n}")
        parts.append("  /  ".join(n_parts))

    axis_label = page.get("axis_label") or cr.get("axis_label", "")
    if axis_label:
        parts.append(f"分析軸: {axis_label}")

    return "    ".join(parts)


def _add_slide(prs: Presentation, slide, page: dict, cr: dict) -> None:
    lc: dict = page.get("layoutConfig") or {}
    cc: dict = page.get("chartConfig") or {}
    mode: str = cc.get("chartMode", "auto")

    title_text = _display_title(page, cr)
    subtitle_text = cr.get("question_text", "") if lc.get("showQuestionText", True) else ""

    title_font_size: int = int(lc.get("titleFontSize", 10))
    footer_font_size: int = int(lc.get("footerFontSize", 4))
    show_footer: bool = lc.get("showFooter", True)

    # タイトル
    _add_text_shape(
        slide, title_text,
        MARGIN_L, MARGIN_T_TITLE, CONTENT_W, TITLE_H,
        font_size=title_font_size, bold=True,
    )
    # サブタイトル（設問文）
    if subtitle_text:
        _add_text_shape(
            slide, subtitle_text,
            MARGIN_L, MARGIN_T_TITLE + TITLE_H, CONTENT_W, SUBTITLE_H,
            font_size=9, bold=False, color="555555",
        )

    is_table_only: bool = (mode == "table_only")

    # tablePosition を解決（後方互換: showTable=True → "bottom"）
    table_pos: str = lc.get("tablePosition", "")
    if not table_pos:
        table_pos = "bottom" if lc.get("showTable", False) else "none"

    show_table: bool = table_pos not in ("none", "separate") or is_table_only

    # フッター分のコンテンツ高さを確保
    footer_reserve = (FOOTER_H + FOOTER_MARGIN_BOTTOM) if show_footer else Inches(0)
    content_h = CONTENT_H - footer_reserve

    GAP = Inches(0.05)
    # 集計表サイズ比率（ユーザー調整可能）
    table_w_pct = lc.get("tableWidthPct", 35) / 100.0
    table_h_pct = lc.get("tableHeightPct", 40) / 100.0

    if is_table_only:
        _add_table_shape(slide, cr, lc, cc, MARGIN_L, CONTENT_TOP, CONTENT_W, content_h)
    elif table_pos == "right" and show_table:
        table_w = CONTENT_W * table_w_pct
        chart_w = CONTENT_W * (1 - table_w_pct) - GAP
        _add_chart_area(slide, cr, cc, mode, MARGIN_L, CONTENT_TOP, chart_w, content_h)
        _add_table_shape(slide, cr, lc, cc, MARGIN_L + chart_w + GAP, CONTENT_TOP, table_w, content_h)
    elif table_pos == "left" and show_table:
        table_w = CONTENT_W * table_w_pct
        chart_w = CONTENT_W * (1 - table_w_pct) - GAP
        _add_table_shape(slide, cr, lc, cc, MARGIN_L, CONTENT_TOP, table_w, content_h)
        _add_chart_area(slide, cr, cc, mode, MARGIN_L + table_w + GAP, CONTENT_TOP, chart_w, content_h)
    elif table_pos == "top" and show_table:
        table_h = content_h * table_h_pct
        chart_h = content_h * (1 - table_h_pct) - GAP
        _add_table_shape(slide, cr, lc, cc, MARGIN_L, CONTENT_TOP, CONTENT_W, table_h)
        _add_chart_area(slide, cr, cc, mode, MARGIN_L, CONTENT_TOP + table_h + GAP, CONTENT_W, chart_h)
    elif show_table:  # bottom
        table_h = content_h * table_h_pct
        chart_h = content_h * (1 - table_h_pct) - GAP
        _add_chart_area(slide, cr, cc, mode, MARGIN_L, CONTENT_TOP, CONTENT_W, chart_h)
        _add_table_shape(slide, cr, lc, cc, MARGIN_L, CONTENT_TOP + chart_h + GAP, CONTENT_W, table_h)
    else:  # none / separate
        _add_chart_area(slide, cr, cc, mode, MARGIN_L, CONTENT_TOP, CONTENT_W, content_h)

    # フッター（n数・分析軸）
    if show_footer:
        footer_text = _build_footer_text(page, cr)
        if footer_text:
            footer_top = CONTENT_TOP + content_h + Inches(0.04)
            _add_text_shape(
                slide, footer_text,
                MARGIN_L, footer_top, CONTENT_W, FOOTER_H,
                font_size=footer_font_size, bold=False, color="555555",
            )


def _add_table_only_slide(prs: Presentation, slide, page: dict, cr: dict) -> None:
    """集計表のみのスライドを生成する（tablePosition="separate" 用）。"""
    lc: dict = page.get("layoutConfig") or {}
    cc: dict = page.get("chartConfig") or {}
    title_text = _display_title(page, cr) + "（集計表）"
    title_font_size: int = int(lc.get("titleFontSize", 18))
    _add_text_shape(
        slide, title_text,
        MARGIN_L, MARGIN_T_TITLE, CONTENT_W, TITLE_H,
        font_size=title_font_size, bold=True,
    )
    _add_table_shape(slide, cr, lc, cc, MARGIN_L, CONTENT_TOP, CONTENT_W, CONTENT_H)


# ---------------------------------------------------------------------------
# グラフ領域
# ---------------------------------------------------------------------------

def _build_split_by_axis_datasets(cr: dict) -> list[dict]:
    """axis_categories ごとに仮想サブデータセットを生成する（by_axis 分割）。"""
    axis_cats   = cr.get("axis_categories") or []
    axis_totals = cr.get("axis_totals") or []
    rows        = cr.get("rows") or []
    choice_labels = [r.get("label", "") for r in rows]
    result = []
    for ci, cat in enumerate(axis_cats):
        percents = [
            r["percents"][ci] if ci < len(r.get("percents", [])) else 0
            for r in rows
        ]
        counts = [
            r["counts"][ci] if isinstance(r.get("counts"), list) and ci < len(r["counts"]) else 0
            for r in rows
        ]
        result.append({
            "target_value":    cat,
            "rows":            [{"label": cat, "percents": percents, "counts": counts}],
            "axis_categories": choice_labels,
            "axis_totals":     [axis_totals[ci] if ci < len(axis_totals) else 0],
        })
    return result


def _build_split_by_comparison_datasets(cr: dict) -> list[dict]:
    """rows ごとに仮想サブデータセットを生成する（by_comparison 分割）。"""
    axis_cats   = cr.get("axis_categories") or []
    axis_totals = cr.get("axis_totals") or []
    rows        = cr.get("rows") or []
    result = []
    for row in rows:
        percents = [
            row["percents"][ci] if ci < len(row.get("percents", [])) else 0
            for ci in range(len(axis_cats))
        ]
        counts = [
            row["counts"][ci] if isinstance(row.get("counts"), list) and ci < len(row["counts"]) else 0
            for ci in range(len(axis_cats))
        ]
        result.append({
            "target_value":    row.get("label", ""),
            "rows":            [{"label": row.get("label", ""), "percents": percents, "counts": counts}],
            "axis_categories": list(axis_cats),
            "axis_totals":     list(axis_totals),
        })
    return result


def _add_chart_area(slide, cr: dict, cc: dict, mode: str, x, y, w, h) -> None:
    """mode / splitMode に応じてグラフを配置する。分割モードは格子状に複数配置。"""
    import math
    split_mode = cc.get("splitMode", "normal")

    # 分割データセットの決定
    if mode == "small_multiples":
        datasets = cr.get("comparison_datasets") or []
    elif split_mode == "by_axis":
        datasets = _build_split_by_axis_datasets(cr)
    elif split_mode == "by_comparison":
        datasets = _build_split_by_comparison_datasets(cr)
    else:
        _add_chart_shape(slide, cr, cc, mode, x, y, w, h)
        return

    if not datasets:
        return

    # ページ分割スライシング: splitDatasetIndices 優先、旧 splitChunkStart/End は fallback
    if split_mode in ("by_axis", "by_comparison"):
        indices = cc.get("splitDatasetIndices")
        if indices is not None:
            datasets = [datasets[i] for i in indices if i < len(datasets)]
        else:
            chunk_start = int(cc.get("splitChunkStart") or 0)
            chunk_end   = int(cc.get("splitChunkEnd")   or len(datasets))
            datasets    = datasets[chunk_start:chunk_end]
        if not datasets:
            return

    # 列数決定 (pageLayout 優先 → splitColumns → 自動)
    n           = len(datasets)
    page_layout = cc.get("pageLayout", "auto")
    if   page_layout in ("cols1", "vertical"):   split_cols = 1
    elif page_layout == "horizontal":            split_cols = n
    elif page_layout in ("cols2", "grid2x2"):   split_cols = 2
    elif page_layout in ("cols3", "grid3x2"):   split_cols = 3
    else:
        split_cols = int(cc.get("splitColumns") or 0) or (1 if n <= 2 else 2 if n <= 4 else 3)

    rows_count = math.ceil(n / split_cols)
    col_w      = w / split_cols
    row_h      = h / rows_count

    for i, ds in enumerate(datasets):
        col = i % split_cols
        row = i // split_cols
        mini_cr = {**cr,
                   "rows":               ds.get("rows", []),
                   "axis_categories":    ds.get("axis_categories", []),
                   "axis_totals":        ds.get("axis_totals", []),
                   "comparison_datasets": None}
        # by_comparison: サブチャート全系列を選択肢の1色で統一
        # STEP3/STEP4 では ds.target_value（選択肢ラベル）の色が全バーに適用される
        if split_mode == "by_comparison":
            choice_label = ds.get("target_value", "")
            choice_color = _resolve_colors(cc, [choice_label])[0] if choice_label else CHART_COLORS[i % len(CHART_COLORS)]
            axis_cats_for_ds = ds.get("axis_categories", [])
            mini_cc = {**cc, "colorSettings": {
                **(cc.get("colorSettings") or {}),
                "overriddenSeriesColors": {cat: choice_color for cat in axis_cats_for_ds},
            }}
        else:
            mini_cc = cc
        _add_chart_shape(slide, mini_cr, mini_cc, "vbar",
                         x + col_w * col, y + row_h * row, col_w, row_h)


def _add_chart_shape(slide, cr: dict, cc: dict, mode: str, x, y, w, h) -> None:
    chart_type = CHART_TYPE_MAP.get(mode, XL_CHART_TYPE.COLUMN_CLUSTERED)
    categories, series_names, data_matrix = _prepare_chart_data(cr, cc, mode)

    if not categories or not series_names:
        return

    # ChartData 組み立て
    chart_data = ChartData()
    chart_data.categories = categories
    for si, sname in enumerate(series_names):
        vals = data_matrix[si] if si < len(data_matrix) else [0] * len(categories)
        chart_data.add_series(sname, vals)

    chart_shape = slide.shapes.add_chart(chart_type, x, y, w, h, chart_data)
    chart = chart_shape.chart

    # 凡例
    if cc.get("showLegend", True):
        from pptx.enum.chart import XL_LEGEND_POSITION
        chart.has_legend = True
        pos_str = cc.get("legendPosition", "bottom")
        pos_map = {
            "bottom": XL_LEGEND_POSITION.BOTTOM,
            "top":    XL_LEGEND_POSITION.TOP,
            "right":  XL_LEGEND_POSITION.RIGHT,
            "left":   XL_LEGEND_POSITION.LEFT,
        }
        chart.legend.position = pos_map.get(pos_str, XL_LEGEND_POSITION.BOTTOM)
        chart.legend.include_in_layout = False
    else:
        chart.has_legend = False

    # 系列カラー
    colors = _resolve_colors(cc, series_names)
    _set_chart_series_colors(chart, colors)

    # データラベル
    if cc.get("showLabels", True):
        from pptx.oxml.ns import qn as pqn
        decimal = int(cc.get("labelDecimalPlaces", 1))
        for plot in chart.plots:
            plot.has_data_labels = True
            dLbls = plot._element.find(pqn("c:dLbls"))
            if dLbls is not None:
                # 数値フォーマット: 0.0% etc.
                numFmt = dLbls.find(pqn("c:numFmt"))
                fmt_code = "0." + "0" * decimal + "%"
                if numFmt is None:
                    numFmt = etree.SubElement(dLbls, pqn("c:numFmt"))
                numFmt.set("formatCode", fmt_code)
                numFmt.set("sourceLinked", "0")


def _prepare_chart_data(
    cr: dict, cc: dict, mode: str
) -> tuple[list[str], list[str], list[list[float]]]:
    """(categories, series_names, data_matrix) を返す。
    data_matrix[series_index][category_index] = float value
    """
    is_brand = mode in ("brand_hbar", "brand_vbar", "brand_vbar_stacked")
    transpose = cc.get("transpose", False)

    if is_brand:
        datasets: list[dict] = cr.get("comparison_datasets") or []
        if not datasets:
            return [], [], []
        categories = [ds.get("target_value", "") for ds in datasets]
        axis_cats: list[str] = (datasets[0].get("axis_categories") or cr.get("axis_categories") or [])
        series_names = axis_cats
        data_matrix: list[list[float]] = []
        for ci in range(len(axis_cats)):
            row_vals: list[float] = []
            for ds in datasets:
                rows = ds.get("rows") or []
                val = rows[0]["percents"][ci] if rows and ci < len(rows[0].get("percents", [])) else 0.0
                row_vals.append(val / 100.0)
            data_matrix.append(row_vals)
    else:
        rows = _filter_rows(cr.get("rows") or [], cc)
        categories = [r["label"] for r in rows]
        axis_cats = cr.get("axis_categories") or []
        series_names = axis_cats
        data_matrix = []
        for ci in range(len(axis_cats)):
            vals = [r["percents"][ci] / 100.0 if ci < len(r.get("percents", [])) else 0.0 for r in rows]
            data_matrix.append(vals)

    if transpose and categories and series_names:
        # 軸を入れ替える
        new_cats = series_names[:]
        new_series = categories[:]
        n_ser = len(new_series)
        n_cat = len(new_cats)
        new_matrix = [[0.0] * n_cat for _ in range(n_ser)]
        for old_si in range(len(data_matrix)):
            for old_ci in range(len(data_matrix[old_si])):
                new_matrix[old_ci][old_si] = data_matrix[old_si][old_ci]
        categories, series_names, data_matrix = new_cats, new_series, new_matrix

    return categories, series_names, data_matrix


# ---------------------------------------------------------------------------
# 集計表
# ---------------------------------------------------------------------------

def _add_table_shape(slide, cr: dict, lc: dict, cc: dict, x, y, w, h) -> None:
    rows_data = _filter_rows(cr.get("rows") or [], cc)
    axis_cats: list[str] = cr.get("axis_categories") or []
    axis_totals: list[int] = cr.get("axis_totals") or []

    content_mode: str = lc.get("tableContentMode", "percent")
    dec_places: int = int(lc.get("tableDecimalPlaces", 1))
    font_size: int = int(lc.get("tableFontSize", 9))
    cell_padding_pt: float | None = lc.get("tableCellPadding")
    # STEP4のlayoutConfig設定を優先、なければSTEP3のchartConfig設定（デフォルトFalse）
    has_total_col: bool = lc.get("showTableRowTotal", cc.get("showTotalCol", False))
    show_col_total: bool = lc.get("showTableColTotal", False)

    # 列定義: ラベル列 → axis_cats → 合計列（右端）
    col_headers = [""] + axis_cats + (["合計"] if has_total_col else [])
    n_cols = len(col_headers)
    axis_col_start = 1                                   # axis_cats の開始列インデックス
    total_col_idx = 1 + len(axis_cats) if has_total_col else None  # 合計列は右端

    # 行定義: ヘッダ行 + N行 + データ行 + (列合計行)
    n_header_rows = 1
    n_n_rows = 1
    data_rows_start = n_header_rows + n_n_rows
    n_data_rows = len(rows_data)
    n_rows_total = n_header_rows + n_n_rows + n_data_rows + (1 if show_col_total else 0)

    table_shape = slide.shapes.add_table(n_rows_total, n_cols, x, y, w, h)
    tbl = table_shape.table

    # 列幅: ラベル列を少し広く、axis列を均等、合計列は少し細め
    label_col_w = int(w * 0.22)
    total_col_w = int(w * 0.08) if has_total_col else 0
    remain_cols  = len(axis_cats)
    axis_col_w   = int((w - label_col_w - total_col_w) / max(remain_cols, 1))
    tbl.columns[0].width = label_col_w
    for ci in range(len(axis_cats)):
        tbl.columns[axis_col_start + ci].width = axis_col_w
    if has_total_col:
        tbl.columns[total_col_idx].width = total_col_w

    # ヘッダ行
    _set_cell(tbl, 0, 0, "", font_size, bold=True)
    for ci, cat in enumerate(axis_cats):
        _set_cell(tbl, 0, axis_col_start + ci, cat, font_size, bold=True)
    if has_total_col:
        _set_cell(tbl, 0, total_col_idx, "合計", font_size, bold=True)

    # N 行
    _set_cell(tbl, 1, 0, "n", font_size, bold=True)
    for ci, total in enumerate(axis_totals):
        _set_cell(tbl, 1, axis_col_start + ci, str(total), font_size, align="center")
    if has_total_col:
        _set_cell(tbl, 1, total_col_idx, str(sum(axis_totals)), font_size, align="center")

    # データ行
    for ri, row in enumerate(rows_data):
        r_idx = data_rows_start + ri
        _set_cell(tbl, r_idx, 0, row.get("label", ""), font_size)
        percents: list[float] = row.get("percents") or []
        counts: list[int] = row.get("counts") or []
        for ci in range(len(axis_cats)):
            pct = percents[ci] if ci < len(percents) else 0.0
            cnt = counts[ci] if ci < len(counts) else 0
            _set_cell(tbl, r_idx, axis_col_start + ci,
                      _format_cell(pct, cnt, content_mode, dec_places),
                      font_size, align="center")
        if has_total_col:
            total_pct = sum(percents[:len(axis_cats)])
            total_cnt = sum(counts[:len(axis_cats)])
            _set_cell(tbl, r_idx, total_col_idx,
                      _format_cell(total_pct, total_cnt, content_mode, dec_places),
                      font_size, align="center")

    # 列合計行（showTableColTotal=True のとき）
    if show_col_total:
        col_total_row_idx = n_header_rows + n_n_rows + n_data_rows
        _set_cell(tbl, col_total_row_idx, 0, "合計", font_size, bold=True)
        for ci in range(len(axis_cats)):
            total_cnt = sum(
                (r.get("counts") or [])[ci] if ci < len(r.get("counts") or []) else 0
                for r in rows_data
            )
            total_n = axis_totals[ci] if ci < len(axis_totals) else 0
            total_pct = total_cnt / total_n * 100 if total_n > 0 else 0.0
            _set_cell(tbl, col_total_row_idx, axis_col_start + ci,
                      _format_cell(total_pct, total_cnt, content_mode, dec_places),
                      font_size, align="center")
        if has_total_col:
            grand_cnt = sum(
                sum((r.get("counts") or [])[:len(axis_cats)]) for r in rows_data
            )
            grand_n = sum(axis_totals)
            grand_pct = grand_cnt / grand_n * 100 if grand_n > 0 else 0.0
            _set_cell(tbl, col_total_row_idx, total_col_idx,
                      _format_cell(grand_pct, grand_cnt, content_mode, dec_places),
                      font_size, align="center")

    # 罫線
    _set_table_borders(tbl, TABLE_BORDER_PT, TABLE_BORDER_COLOR)

    # セル内余白（指定がある場合のみ上書き）
    if cell_padding_pt is not None:
        _set_table_cell_margins(tbl, float(cell_padding_pt))


def _format_cell(pct: float, cnt: int, mode: str, dec: int) -> str:
    fmt = f"{{:.{dec}f}}%"
    if mode == "count":
        return str(cnt)
    elif mode == "both":
        return f"{fmt.format(pct)}\n({cnt})"
    return fmt.format(pct)


def _set_cell(tbl, row: int, col: int, text: str, font_size: int,
              bold: bool = False, align: str = "left") -> None:
    from pptx.enum.text import PP_ALIGN
    cell = tbl.cell(row, col)
    cell.text = text
    tf = cell.text_frame
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER if align == "center" else PP_ALIGN.LEFT
        for run in para.runs:
            run.font.name = FONT_NAME
            run.font.size = Pt(font_size)
            run.font.bold = bold


def _set_table_cell_margins(tbl, pad_pt: float) -> None:
    """全セルの内部余白を pt 単位で設定する（XML操作）。"""
    from pptx.oxml.ns import qn as pqn
    v_emu = int(Pt(pad_pt))
    h_emu = int(Pt(pad_pt * 2.5))
    for row_idx in range(len(tbl.rows)):
        for col_idx in range(len(tbl.columns)):
            cell = tbl.cell(row_idx, col_idx)
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            tcPr.set("marL", str(h_emu))
            tcPr.set("marR", str(h_emu))
            tcPr.set("marT", str(v_emu))
            tcPr.set("marB", str(v_emu))


def _set_table_borders(tbl, border_pt: float, color_hex: str) -> None:
    """XML 操作で全セルの罫線を設定する。"""
    emu = int(border_pt * 12700)  # 1pt = 12700 EMU
    nsmap = "http://schemas.openxmlformats.org/drawingml/2006/main"

    def _make_ln(tag_name: str) -> etree._Element:
        ln = etree.Element(f"{{{nsmap}}}{tag_name}")
        ln.set("w", str(emu))
        solidFill = etree.SubElement(ln, f"{{{nsmap}}}solidFill")
        srgbClr = etree.SubElement(solidFill, f"{{{nsmap}}}srgbClr")
        srgbClr.set("val", color_hex)
        prstDash = etree.SubElement(ln, f"{{{nsmap}}}prstDash")
        prstDash.set("val", "solid")
        return ln

    for row_idx in range(len(tbl.rows)):
        for col_idx in range(len(tbl.columns)):
            tc = tbl.cell(row_idx, col_idx)._tc
            tcPr = tc.find(f"{{{nsmap}}}tcPr")
            if tcPr is None:
                tcPr = etree.SubElement(tc, f"{{{nsmap}}}tcPr")
            for side in ("lnL", "lnR", "lnT", "lnB"):
                existing = tcPr.find(f"{{{nsmap}}}{side}")
                if existing is not None:
                    tcPr.remove(existing)
                tcPr.append(_make_ln(side))


# ---------------------------------------------------------------------------
# テキストボックス
# ---------------------------------------------------------------------------

def _add_text_shape(slide, text: str, x, y, w, h,
                    font_size: int = 12, bold: bool = False,
                    color: str = "000000") -> None:
    from pptx.enum.text import PP_ALIGN
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


# ---------------------------------------------------------------------------
# カラー解決（JS の _resolveColorsForPage と同一ロジック）
# ---------------------------------------------------------------------------

def _resolve_colors(cc: dict, labels: list[str]) -> list[str]:
    cs: dict = cc.get("colorSettings") or {}
    # JSはArray<{label,color}>形式なのでdictに変換
    vcm_list = cs.get("valueColorMapping") or []
    value_color_map: dict = {
        e["label"]: e["color"]
        for e in (vcm_list if isinstance(vcm_list, list) else [])
        if isinstance(e, dict) and "label" in e
    }
    overridden: dict = cs.get("overriddenSeriesColors") or {}
    resolved_color_map: dict = cs.get("resolvedColorMap") or {}
    # selectedPalette キーが存在してかつ None = 明示的グレーパレット
    is_gray_palette = "selectedPalette" in cs and cs.get("selectedPalette") is None

    colors: list[str] = []
    palette_idx = 0
    for label in labels:
        if label in overridden:                                    # 手動上書き（最優先）
            colors.append(overridden[label])
        elif label in resolved_color_map:                          # STEP3確定色
            colors.append(resolved_color_map[label])
        elif label == "その他":
            colors.append("#aaaaaa")
        elif label == "全体":
            colors.append("#555555")
        elif label in value_color_map:                             # valueColorMapping
            colors.append(value_color_map[label])
        elif is_gray_palette:                                      # グレーパレット
            colors.append("#767676")
        else:
            colors.append(CHART_COLORS[palette_idx % len(CHART_COLORS)])
            palette_idx += 1
    return colors


def _set_chart_series_colors(chart, colors: list[str]) -> None:
    from pptx.dml.color import RGBColor as RC
    for plot in chart.plots:
        for si, series in enumerate(plot.series):
            color_hex = colors[si % len(colors)] if colors else "#0071BC"
            color_hex = color_hex.lstrip("#")
            pt = series.format.fill
            pt.solid()
            pt.fore_color.rgb = RC.from_string(color_hex)


# ---------------------------------------------------------------------------
# 行フィルタ・並び替え
# ---------------------------------------------------------------------------

def _filter_rows(rows: list[dict], cc: dict) -> list[dict]:
    hidden: set[str] = set(cc.get("hiddenChoices") or [])
    row_order: list[str] | None = (cc.get("layoutConfig") or {}).get("rowChoiceOrder") or cc.get("rowChoiceOrder")
    sort_order: str = cc.get("sortOrder") or "original"

    filtered = [r for r in rows if r.get("label") not in hidden]

    if row_order and (not sort_order or sort_order == "original"):
        order_map = {lbl: i for i, lbl in enumerate(row_order)}
        in_order = sorted([r for r in filtered if r.get("label") in order_map],
                          key=lambda r: order_map[r["label"]])
        rest = [r for r in filtered if r.get("label") not in order_map]
        filtered = in_order + rest

    if sort_order == "desc":
        # 最初の系列の percents[0] でソート
        filtered.sort(key=lambda r: -(r.get("percents") or [0])[0])
    elif sort_order == "asc":
        filtered.sort(key=lambda r: (r.get("percents") or [0])[0])

    return filtered


# ---------------------------------------------------------------------------
# タイトル表示文字列
# ---------------------------------------------------------------------------

def _display_title(page: dict, cr: dict) -> str:
    lc: dict = page.get("layoutConfig") or {}
    override = lc.get("titleOverride") or None
    return override if override else (cr.get("title") or "")
